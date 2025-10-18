from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.util import Cm

LOGO_PATH = "logo.png"  # Remplace par ton fichier
FONT_NAME = "Calibri"
TITLE_COLOR = RGBColor(0, 70, 127)  # Bleu pro
TEXT_COLOR = RGBColor(0, 0, 0)

def add_title_slide(prs, title, subtitle):
    slide_layout = prs.slide_layouts[6]  # vide
    slide = prs.slides.add_slide(slide_layout)

    # Ajout titre
    title_shape = slide.shapes.add_textbox(Cm(2), Cm(5), Cm(24), Cm(3))
    tf = title_shape.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(32)
    run.font.bold = True
    run.font.color.rgb = TITLE_COLOR
    run.font.name = FONT_NAME
    p.alignment = PP_ALIGN.LEFT

    # Ajout sous-titre
    subtitle_shape = slide.shapes.add_textbox(Cm(2), Cm(8), Cm(24), Cm(3))
    tf2 = subtitle_shape.text_frame
    p2 = tf2.paragraphs[0]
    run2 = p2.add_run()
    run2.text = subtitle
    run2.font.size = Pt(20)
    run2.font.name = FONT_NAME
    p2.alignment = PP_ALIGN.LEFT

    # Ajout logo (en haut à droite)
    slide.shapes.add_picture(LOGO_PATH, Cm(22), Cm(1), height=Cm(2.5))

def add_content_slide(prs, title, bullet_points):
    slide_layout = prs.slide_layouts[6]  # vide
    slide = prs.slides.add_slide(slide_layout)

    # Titre
    title_shape = slide.shapes.add_textbox(Cm(1), Cm(1), Cm(24), Cm(2))
    tf = title_shape.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = TITLE_COLOR
    run.font.name = FONT_NAME

    # Contenu
    body_shape = slide.shapes.add_textbox(Cm(2), Cm(4), Cm(24), Cm(15))
    tf_body = body_shape.text_frame
    for point in bullet_points:
        p = tf_body.add_paragraph()
        p.text = point
        p.level = 0
        p.font.size = Pt(18)
        p.font.name = FONT_NAME
        p.font.color.rgb = TEXT_COLOR

def main():
    prs = Presentation()

    add_title_slide(prs,
        "Parcours Compétences & Retour à l’Emploi",
        "Présentation du projet\nDate : 2025"
    )

    # Slides suivantes (résumé)
    slides = [
        ("Contexte et Enjeux", [
            "Accompagner les demandeurs d’emploi vers un retour durable à l’emploi",
            "Évaluation personnalisée des compétences",
            "Orientation vers les métiers porteurs"
        ]),
        ("Objectifs", [
            "Identifier compétences techniques et transversales",
            "Faire émerger des pistes métiers réalistes",
            "Créer un plan d’action personnalisé"
        ]),
        ("Public Cible", [
            "Demandeurs d’emploi (France Travail)",
            "DE longue durée, jeunes sans qualification, seniors, reconversion"
        ]),
        ("Format et Déroulé", [
            "Atelier collectif + rendez-vous individuel",
            "3 à 4 demi-journées",
            "Fréquence : mensuelle ou bimensuelle"
        ]),
        ("Contenu par Séance", [
            "1. Diagnostic & Profil de compétences",
            "2. Métiers et secteurs compatibles",
            "3. Stratégie de recherche d’emploi",
            "4. Préparation aux entretiens (optionnel)"
        ]),
        ("Livrables", [
            "Fiche profil DE : compétences + métiers + plan d’action",
            "CV actualisé + lettre personnalisée",
            "Suivi post-atelier"
        ]),
        ("Partenaires mobilisés", [
            "France Travail (entreprises / formations)",
            "Structures insertion / PLIE / Missions locales",
            "Diagoriente, Emploi Store, La Bonne Boîte"
        ]),
        ("Indicateurs de réussite", [
            "% DE avec une orientation métiers",
            "% DE ayant candidaté sous 30 jours",
            "% DE ayant obtenu un entretien ou une formation",
            "Taux de satisfaction global"
        ]),
        ("Budget prévisionnel", [
            "Temps animateur / conseiller",
            "Impression des supports",
            "Outils numériques (licences éventuelles)"
        ]),
        ("Prochaines étapes", [
            "Finalisation des supports",
            "Lancement pilote",
            "Évaluation et déploiement progressif"
        ])
    ]

    for title, points in slides:
        add_content_slide(prs, title, points)

    prs.save("Projet_Parcours_Competences_Retour_Emploi_CUSTOM.pptx")
    print("✅ Présentation PPTX personnalisée générée avec succès !")

if __name__ == "__main__":
    main()
